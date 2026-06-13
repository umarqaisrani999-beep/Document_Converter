import os
import fitz  
from pptx import Presentation 
from pptx.util import Inches

def pdf_to_pptx(input_path: str, output_path: str, progress_cb=None):
    
    if progress_cb:
        progress_cb(5, "Opening PDF document...")

    
    doc = fitz.open(input_path)
    total_pages = len(doc)
    
    if total_pages == 0:
        raise ValueError("The selected PDF file contains no pages.")

    
    prs = Presentation()
    
    
    blank_slide_layout = prs.slide_layouts[6] 

    if progress_cb:
        progress_cb(15, f"Converting {total_pages} pages...")

   
    temp_images = []

    try:
        for page_num in range(total_pages):
            page = doc[page_num]
            
            
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            
            
            image_filename = f"temp_page_{page_num}.png"
            pix.save(image_filename)
            temp_images.append(image_filename)

            
            if page_num == 0:
                prs.slide_width = Inches(pix.width / 96)
                prs.slide_height = Inches(pix.height / 96)

            
            slide = prs.slides.add_slide(blank_slide_layout)

            
            slide.shapes.add_picture(image_filename, 0, 0, width=prs.slide_width, height=prs.slide_height)

            
            if progress_cb:
                pct = 15 + int((page_num + 1) / total_pages * 75)
                progress_cb(min(pct, 90), f"Processing page {page_num + 1}/{total_pages}...")

        if progress_cb:
            progress_cb(95, "Saving PowerPoint file...")
            
        
        prs.save(output_path)

        if progress_cb:
            progress_cb(100, "Conversion complete!")

    finally:
        
        doc.close()
        
        for img_path in temp_images:
            if os.path.exists(img_path):
                try:
                    os.remove(img_path)
                except Exception:
                    pass  