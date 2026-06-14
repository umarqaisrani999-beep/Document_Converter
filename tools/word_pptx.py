import os
from docx import Document
from pptx import Presentation


def word_to_pptx(input_path: str, output_path: str, progress_cb=None):
    
    if progress_cb:
        progress_cb(5, "Opening Word document...")

    doc = Document(input_path)

    paragraphs = [p for p in doc.paragraphs if p.text.strip()]
    total = len(paragraphs)

    if total == 0:
        raise ValueError("The selected Word document has no readable text.")

    if progress_cb:
        progress_cb(15, "Reading document structure...")

    # --- Step 1: turn the flat list of paragraphs into a slide outline ---
    # Each slide is: {"title": str, "bullets": [(text, is_subheading), ...]}
    MAX_BULLETS_PER_SLIDE = 6
    slides_data = []
    current_slide = None

    for i, para in enumerate(paragraphs):
        text = para.text.strip()
        style_name = (para.style.name if para.style else "") or ""

        is_heading = style_name.startswith("Heading") or style_name == "Title"
        heading_level = 0
        if style_name == "Title":
            heading_level = 1
        elif style_name.startswith("Heading"):
            try:
                heading_level = int(style_name.replace("Heading", "").strip())
            except ValueError:
                heading_level = 1

        if is_heading and heading_level == 1:
            # Top-level heading -> start a new slide
            current_slide = {"title": text, "bullets": []}
            slides_data.append(current_slide)
        else:
            if current_slide is None:
                # No top-level heading seen yet -> create a default slide
                current_slide = {"title": "Untitled Slide", "bullets": []}
                slides_data.append(current_slide)

            current_slide["bullets"].append((text, is_heading))

            # Avoid cramming too much text onto a single slide
            if len(current_slide["bullets"]) >= MAX_BULLETS_PER_SLIDE:
                current_slide = {"title": current_slide["title"] + " (cont.)", "bullets": []}
                slides_data.append(current_slide)

        if progress_cb:
            pct = 15 + int((i + 1) / total * 35)
            progress_cb(min(pct, 50), "Reading document content...")

    # Drop any empty trailing "(cont.)" slide created by the loop above
    if slides_data and not slides_data[-1]["bullets"] and len(slides_data) > 1:
        slides_data.pop()

    if progress_cb:
        progress_cb(55, "Building presentation...")

    # --- Step 2: build the actual PPTX from the outline ---
    prs = Presentation()
    content_layout = prs.slide_layouts[1]  # "Title and Content"

    total_slides = len(slides_data)

    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(content_layout)

        if slide.shapes.title is not None:
            slide.shapes.title.text = slide_data["title"] or f"Slide {idx + 1}"

        if slide_data["bullets"]:
            body = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx != 0:  # idx 0 is always the title
                    body = shape
                    break

            if body is not None:
                tf = body.text_frame
                tf.clear()
                first = True
                for text, is_subheading in slide_data["bullets"]:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.text = text
                    # Sub-headings sit at the top level; regular text is indented
                    p.level = 0 if is_subheading else 1
                    if is_subheading:
                        for run in p.runs:
                            run.font.bold = True

        if progress_cb:
            pct = 55 + int((idx + 1) / total_slides * 40)
            progress_cb(min(pct, 95), f"Creating slide {idx + 1}/{total_slides}...")

    if progress_cb:
        progress_cb(98, "Saving PowerPoint file...")

    prs.save(output_path)

    if not os.path.exists(output_path):
        raise RuntimeError("Word to PPTX conversion did not produce an output file.")
 
    if progress_cb:
        progress_cb(100, "Conversion complete!")
 
    return output_path