import os
from pptx import Presentation
from docx import Document


def pptx_to_word(input_path: str, output_path: str, progress_cb=None):
    
    if progress_cb:
        progress_cb(5, "Opening PowerPoint file...")

    prs = Presentation(input_path)
    total_slides = len(prs.slides)

    if total_slides == 0:
        raise ValueError("The selected PPTX file contains no slides.")

    if progress_cb:
        progress_cb(15, f"Processing {total_slides} slide(s)...")

    doc = Document()
    doc.add_heading("Presentation Content", level=0)

    for idx, slide in enumerate(prs.slides):
        title_text = None
        body_blocks = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            # The title placeholder always has placeholder idx == 0.
            # (Comparing shape objects directly via `is` is unreliable in
            # python-pptx, so we check the placeholder index instead.)
            is_title_placeholder = False
            if shape.is_placeholder:
                try:
                    is_title_placeholder = shape.placeholder_format.idx == 0
                except Exception:
                    is_title_placeholder = False

            if is_title_placeholder and title_text is None:
                title_text = text
            else:
                body_blocks.append(text)

        heading_text = title_text if title_text else f"Slide {idx + 1}"
        doc.add_heading(heading_text, level=1)

        for block in body_blocks:
            for line in block.split("\n"):
                line = line.strip()
                if line:
                    doc.add_paragraph(line, style="List Bullet")

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                p = doc.add_paragraph()
                run = p.add_run(f"Notes: {notes_text}")
                run.italic = True

        if idx < total_slides - 1:
            doc.add_page_break()

        if progress_cb:
            pct = 15 + int((idx + 1) / total_slides * 75)
            progress_cb(min(pct, 90), f"Processing slide {idx + 1}/{total_slides}...")

    if progress_cb:
        progress_cb(95, "Saving Word document...")

    doc.save(output_path)

    if not os.path.exists(output_path):
        raise RuntimeError("PPTX to Word conversion did not produce an output file.")

    if progress_cb:
        progress_cb(100, "Conversion complete!")

    return output_path